#!/usr/bin/env python3
"""Build an EDITABLE A1 .pptx (real text boxes + shapes) for Canva/PowerPoint.

Unlike build_pptx.py (which wraps the flat PNG), this rebuilds the poster as
native, editable elements. The comic effects, custom Google fonts and
animations from poster.html do NOT carry over — Canva/PowerPoint will use their
own fonts — but every text and colour is editable. The UTM/DIGITEX frame is kept
as a background image, and the stylised phone is dropped in as phone.png.

Run first:  node render.mjs && node render_phone.mjs
Then:       pip install python-pptx && python build_editable_pptx.py
"""
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
FRAME = HERE / "template-frame.png"
QR = HERE / "qr.png"
PHONE = HERE / "phone.png"
OUT = HERE / "MentorMinds-DIGITEX2026-editable.pptx"

# Shapes queued for PowerPoint entrance animations (editable in Present mode).
ANIM_TARGETS: list[tuple] = []

# Design canvas is 2245x3179 px (96 dpi) == A1 portrait. 1px = 9525 EMU.
EMU = 9525
def px(v): return Emu(int(round(v * EMU)))
def fs(v):  # design px font-size -> points
    return Pt(round(v * 0.75, 1))

C = {
    "magenta": RGBColor(0xFF, 0x5C, 0xA3),
    "purple":  RGBColor(0x7C, 0x3A, 0xED),
    "ink":     RGBColor(0x14, 0x0F, 0x1A),
    "paper":   RGBColor(0xFF, 0xF5, 0xFA),
    "muted":   RGBColor(0x3A, 0x2F, 0x42),
    "white":   RGBColor(0xFF, 0xFF, 0xFF),
    "yellow":  RGBColor(0xFF, 0xD2, 0x3F),
    "cyan":    RGBColor(0x16, 0xA5, 0xC4),
    "sdg4":    RGBColor(0xC5, 0x19, 0x2D),
    "sdg10":   RGBColor(0xDD, 0x13, 0x67),
}
FONT = "Poppins"
MONO = "Roboto Mono"

prs = Presentation()
prs.slide_width = px(2245)
prs.slide_height = px(3179)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
shapes = slide.shapes

# Official UTM / DIGITEX frame as the background.
shapes.add_picture(str(FRAME), 0, 0, width=px(2245), height=px(3179))


def box(left, top, w, h, fill=None, line=None, line_px=0,
        kind=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06):
    sp = shapes.add_shape(kind, px(left), px(top), px(w), px(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if not line or not line_px:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(round(line_px * 0.75, 1))
    sp.shadow.inherit = False
    if kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def text(left, top, w, h, paragraphs, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.06, wrap=True, animate=False):
    """paragraphs: list of paragraphs; each paragraph is a list of run dicts
    {t, size, color, bold, italic, font}."""
    tb = shapes.add_textbox(px(left), px(top), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for run in para:
            r = p.add_run(); r.text = run["t"]
            f = r.font
            f.size = fs(run.get("size", 24))
            f.bold = run.get("bold", True)
            f.italic = run.get("italic", False)
            f.name = run.get("font", FONT)
            f.color.rgb = run.get("color", C["ink"])
    if animate:
        ANIM_TARGETS.append(tb)
    return tb


def apply_fade_animations(slide, targets, duration_ms=450, stagger_ms=180):
    """Add staggered fade-in entrance animations (PowerPoint Present mode)."""
    if not targets:
        return

    root = slide.element
    for old in root.findall(qn("p:timing")):
        root.remove(old)

    P = qn("p:")
    timing = etree.SubElement(root, P + "timing")
    tnLst = etree.SubElement(timing, P + "tnLst")
    par = etree.SubElement(tnLst, P + "par")
    cTn_root = etree.SubElement(
        par, P + "cTn", id="1", dur="indefinite", restart="never", nodeType="tmRoot"
    )
    childTnLst = etree.SubElement(cTn_root, P + "childTnLst")
    seq = etree.SubElement(childTnLst, P + "seq", concurrent="1", nextAc="seek")
    cTn_seq = etree.SubElement(
        seq, P + "cTn", id="2", dur="indefinite", nodeType="mainSeq"
    )
    childTnLst2 = etree.SubElement(cTn_seq, P + "childTnLst")

    nid = 3
    for idx, shape in enumerate(targets):
        delay = idx * stagger_ms
        spid = str(shape.shape_id)

        par_anim = etree.SubElement(childTnLst2, P + "par")
        cTn_par = etree.SubElement(par_anim, P + "cTn", id=str(nid), fill="hold")
        stCondLst = etree.SubElement(cTn_par, P + "stCondLst")
        etree.SubElement(stCondLst, P + "cond", delay=str(delay))
        childTnLst3 = etree.SubElement(cTn_par, P + "childTnLst")

        par_eff = etree.SubElement(childTnLst3, P + "par")
        cTn_eff = etree.SubElement(
            par_eff, P + "cTn",
            id=str(nid + 1), presetClass="entr", presetID="10",
            fill="hold", nodeType="clickEffect" if idx == 0 else "withEffect",
        )
        stCondLst2 = etree.SubElement(cTn_eff, P + "stCondLst")
        etree.SubElement(stCondLst2, P + "cond", delay="0")
        childTnLst4 = etree.SubElement(cTn_eff, P + "childTnLst")

        par_set = etree.SubElement(childTnLst4, P + "par")
        cTn_set = etree.SubElement(
            par_set, P + "cTn", id=str(nid + 2), fill="hold", presetClass="entr", presetID="10"
        )
        stCondLst3 = etree.SubElement(cTn_set, P + "stCondLst")
        etree.SubElement(stCondLst3, P + "cond", delay="0")
        childTnLst5 = etree.SubElement(cTn_set, P + "childTnLst")

        etree.SubElement(childTnLst5, P + "set", manual="1")
        cBhvr = etree.SubElement(childTnLst5, P + "animEffect", transition="in", filter="fade")
        cBhvr2 = etree.SubElement(cBhvr, P + "cBhvr")
        etree.SubElement(cBhvr2, P + "cTn", id=str(nid + 3), dur=str(duration_ms), fill="hold")
        tgtEl = etree.SubElement(cBhvr2, P + "tgtEl")
        etree.SubElement(tgtEl, P + "spTgt", spid=spid)

        nid += 10


def heading(panel_left, panel_top, num, label, pw):
    d = 56
    box(panel_left + 30, panel_top + 26, d, d, fill=C["magenta"],
        line=C["ink"], line_px=4, kind=MSO_SHAPE.OVAL)
    text(panel_left + 30, panel_top + 26, d, d,
         [[{"t": str(num), "size": 33, "color": C["white"]}]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(panel_left + 30 + d + 18, panel_top + 24, pw - 30 - d - 48, 60,
         [[{"t": label, "size": 39, "color": C["magenta"]}]],
         anchor=MSO_ANCHOR.MIDDLE)


# ---------------- header ----------------
text(132, 466, 1500, 40,
     [[{"t": "AI-POWERED PERSONAL TUTOR · GROUNDED IN REAL EXAM MARK SCHEMES",
        "size": 27, "color": C["magenta"]}]])
text(122, 500, 1500, 180,
     [[{"t": "Mentor", "size": 150, "color": C["magenta"]},
       {"t": "Minds", "size": 150, "color": C["purple"]}]], line_spacing=0.9)
text(134, 678, 1500, 110,
     [[{"t": "An examiner-grade AI tutor for every learner — available 24/7.",
        "size": 42, "color": C["ink"]}]])

# QR card
box(1773, 468, 344, 332, fill=C["white"], line=C["ink"], line_px=6, radius=0.08)
text(1773, 490, 344, 36,
     [[{"t": "Scan for live demo & source", "size": 24, "color": C["magenta"]}]],
     align=PP_ALIGN.CENTER)
shapes.add_picture(str(QR), px(1839), px(534), width=px(212), height=px(212))
text(1773, 752, 344, 30,
     [[{"t": "github.com/", "size": 18, "color": C["ink"], "font": MONO},
       {"t": "Arnobrizwan/mentormind", "size": 18, "color": C["purple"], "font": MONO}]],
     align=PP_ALIGN.CENTER)

# ---------------- 1 · Problem ----------------
box(110, 828, 2025, 352, fill=C["white"], line=C["ink"], line_px=6)
heading(110, 828, 1, "The Problem", 2025)
text(140, 922, 1965, 70,
     [[{"t": "Most students cannot get help the moment they are stuck. Quality "
            "tutoring remains expensive, scarce, and one-size-fits-all — so "
            "doubts pile up and learning gaps widen.",
        "size": 25, "color": C["muted"], "bold": False}]])

problems = [
    ("🎓  No personalized tutoring",
     "One teacher to 40+ students means individual questions go unanswered, and weaker learners fall behind."),
    ("💸  High tuition costs",
     "Private tutoring runs RM50–150 per hour — out of reach for the families who need support the most."),
    ("📡  Limited access to support",
     "Rural and low-income learners lack reliable, on-demand help after school hours or offline."),
]
cw = 637
for i, (h, body) in enumerate(problems):
    l = 140 + i * (cw + 26)
    box(l, 1000, cw, 152, fill=C["paper"], line=C["ink"], line_px=4, radius=0.08)
    text(l + 22, 1016, cw - 44, 40, [[{"t": h, "size": 26, "color": C["ink"]}]])
    text(l + 22, 1062, cw - 44, 84,
         [[{"t": body, "size": 22, "color": C["muted"], "bold": False}]])

# ---------------- 2 · Why It's Different ----------------
box(110, 1210, 660, 1008, fill=C["white"], line=C["ink"], line_px=6)
heading(110, 1210, 2, "Why It's Different", 660)
box(140, 1305, 600, 218, fill=C["magenta"], line=C["ink"], line_px=4, radius=0.1)
text(162, 1322, 556, 188,
     [[{"t": "NOVEL CONTRIBUTION", "size": 18, "color": C["yellow"], "font": MONO}],
      [{"t": "A tutor that grounds every answer in real Cambridge mark schemes "
            "— a self-hosted retrieval model returns examiner working and cites "
            "the exact paper, with an optional fine-tuned (LoRA) model. "
            "No third-party AI APIs.",
        "size": 24, "color": C["white"], "bold": True}]], line_spacing=1.12)

vs = [
    ("vs. general chatbots (ChatGPT)",
     "Replies are drawn from the official mark scheme and cite their source — not hallucinated — and cost $0 in AI API fees."),
    ("vs. Khan Academy",
     "Built around the exact Cambridge syllabus (20 subjects) with examiner-grade worked solutions, on a fully self-hosted, free-tier stack."),
    ("vs. Duolingo",
     "Covers full academic subjects — Maths & Sciences — with step-by-step solutions, not just language drills."),
]
y = 1305 + 218 + 12  # novel box + tight gap
for tag, body in vs:
    row_h = 108
    text(140, y, 600, row_h,
         [[{"t": tag, "size": 23, "color": C["purple"]}],
          [{"t": body, "size": 24, "color": C["muted"], "bold": False}]],
         line_spacing=1.08, animate=True)
    y += row_h + 6

# ---------------- phone hero (image) ----------------
shapes.add_picture(str(PHONE), px(838), px(1220), width=px(585), height=px(1018))
text(850, 2228, 545, 30,
     [[{"t": "↑ The live app: a mark-scheme-grounded solution, a quiz, and progress tracking",
        "size": 20, "color": C["purple"], "font": MONO}]], align=PP_ALIGN.CENTER)

# ---------------- 3 · AI Under the Hood ----------------
box(1475, 1210, 660, 1008, fill=C["white"], line=C["ink"], line_px=6)
heading(1475, 1210, 3, "AI Under the Hood", 660)
tech = [
    ("🧠", "Machine Learning",
     "A self-hosted retrieval model matches questions to the mark-scheme corpus, with an optional LoRA fine-tune (Qwen2.5) — no third-party AI APIs."),
    ("💬", "Natural Language Processing",
     "Understands a typed question, finds the closest Cambridge mark scheme, and explains the answer clearly, step by step."),
    ("📷", "Computer Vision",
     "OpenCV powers exam proctoring (face detection), OMR bubble-sheet auto-grading, and Tesseract OCR for scanned work."),
    ("📈", "Predictive Analytics & MLOps",
     "A dropout-risk model flags at-risk learners and recommends courses, shipped through a DVC + MLflow pipeline with a drift gate."),
]
icon_bg = [C["paper"], C["paper"], C["paper"], C["paper"]]
y = 1300
for i, (ico, h, body) in enumerate(tech):
    row_h = 132
    box(1505, y, 70, 70, fill=icon_bg[i], line=C["ink"], line_px=4, radius=0.18)
    text(1505, y, 70, 70, [[{"t": ico, "size": 36}]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, animate=True)
    text(1593, y, 512, row_h,
         [[{"t": h, "size": 26, "color": C["ink"]}],
          [{"t": body, "size": 22, "color": C["muted"], "bold": False}]],
         line_spacing=1.06, animate=True)
    y += row_h + 8

# ---------------- 4 · Evidence & Impact ----------------
box(110, 2252, 1230, 368, fill=C["white"], line=C["ink"], line_px=6)
heading(110, 2252, 4, "Evidence & Impact", 1230)
metrics = [("1,142", "mark-scheme answers", C["magenta"]),
           ("20", "Cambridge subjects", C["purple"]),
           ("$0", "third-party AI cost", C["cyan"])]
mw = 374
for i, (big, cap, col) in enumerate(metrics):
    l = 140 + i * (mw + 24)
    box(l, 2356, mw, 116, fill=C["paper"], line=C["ink"], line_px=4, radius=0.1)
    text(l, 2366, mw, 76, [[{"t": big, "size": 64, "color": col}]],
         align=PP_ALIGN.CENTER)
    text(l, 2440, mw, 28, [[{"t": cap, "size": 20, "color": C["ink"]}]],
         align=PP_ALIGN.CENTER)
text(140, 2490, 1170, 120,
     [[{"t": "Evidence: ", "size": 21, "color": C["purple"]},
       {"t": "1,521 Cambridge past-paper documents ingested · every answer cites "
            "its exact mark scheme · k6 load-tested · 100% self-hosted on a "
            "free-tier stack, available 24/7.", "size": 21, "color": C["muted"], "bold": False}],
      [{"t": "Classroom pilot in progress — learning-gain & satisfaction results to follow.",
        "size": 16, "color": C["muted"], "bold": False, "italic": True, "font": MONO}]],
     line_spacing=1.15)

# ---------------- 5 · UN SDG Alignment ----------------
box(1380, 2252, 755, 368, fill=C["white"], line=C["ink"], line_px=6)
heading(1380, 2252, 5, "UN SDG Alignment", 755)
sdg = [("4", "Quality Education",
        "Free, examiner-grade tutoring for every learner, any time.", C["sdg4"]),
       ("10", "Reduced Inequalities",
        "The same support for rural & low-income students as elite tutoring.", C["sdg10"])]
sw = 336
for i, (num, goal, body, col) in enumerate(sdg):
    l = 1410 + i * (sw + 22)
    box(l, 2356, sw, 232, fill=col, line=C["ink"], line_px=5, radius=0.08)
    text(l + 20, 2368, sw - 40, 80, [[{"t": num, "size": 60, "color": C["white"]}]])
    text(l + 20, 2452, sw - 40, 40, [[{"t": goal, "size": 23, "color": C["white"]}]])
    text(l + 20, 2496, sw - 40, 80,
         [[{"t": body, "size": 20, "color": C["white"], "bold": False}]], line_spacing=1.08)

# ---------------- footer ----------------
text(110, 2660, 2025, 60,
     [[{"t": "Full-stack: Django · FastAPI · Angular · self-hosted retrieval + "
            "LoRA tutor · OpenCV · Tesseract · Mistral-OCR ingest · DVC + MLflow "
            "· Redis · Postgres · Docker · Kubernetes",
        "size": 26, "color": C["ink"], "bold": False, "font": MONO}]],
     align=PP_ALIGN.CENTER, line_spacing=1.1)
text(110, 2724, 2025, 40,
     [[{"t": "Arnob Rizwan Ahmad · A22EC4005 · Faculty of Computing, UTM",
        "size": 31, "color": C["ink"]}]], align=PP_ALIGN.CENTER)
text(110, 2770, 2025, 36,
     [[{"t": "Mentor: Dr. Ruhaidah Binti Samsudin", "size": 27, "color": C["purple"]}]],
     align=PP_ALIGN.CENTER)

apply_fade_animations(slide, ANIM_TARGETS, duration_ms=450, stagger_ms=160)

prs.save(str(OUT))
print(f"wrote {OUT.name} ({len(ANIM_TARGETS)} animated elements)")
