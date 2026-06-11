#!/usr/bin/env python3
"""Wrap the rendered poster PNG into an A1-portrait .pptx for print/submission.

Run `node render.mjs` first to refresh the PNG, then:
    pip install python-pptx
    python build_pptx.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).resolve().parent
PNG = HERE / "MentorMinds-DIGITEX2026-poster.png"
OUT = HERE / "MentorMinds-DIGITEX2026-poster.pptx"

# A1 portrait: 594mm x 841mm (1mm = 36000 EMU)
MM = 36000
SLIDE_W = Emu(594 * MM)
SLIDE_H = Emu(841 * MM)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
slide.shapes.add_picture(str(PNG), 0, 0, width=SLIDE_W, height=SLIDE_H)
prs.save(str(OUT))
print(f"wrote {OUT.name}")
